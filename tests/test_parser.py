"""문서 파서 (FR-02, 04) 단위 테스트 — 픽스처는 테스트 안에서 생성한다."""

import io

from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter

from app.services.parser import parse_document


def _make_pptx() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 제목만 있는 레이아웃
    slide.shapes.title.text = "금 시세 조회 자동화"

    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "네이버 증권에서 국내 금 시세를 조회한다"

    table = slide.shapes.add_table(2, 2, Inches(1), Inches(4), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "단계"
    table.cell(0, 1).text = "시스템"
    table.cell(1, 0).text = "시세 조회"
    table.cell(1, 1).text = "Edge"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_blank_pdf() -> bytes:
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def test_pptx_extracts_text_and_table_in_order():
    result = parse_document("sample.pptx", _make_pptx())

    assert result["parser"] == "python-pptx"
    assert result["page_count"] == 1
    blocks = result["pages"][0]["blocks"]
    types = [b["type"] for b in blocks]
    assert "text" in types and "table" in types

    # 읽기 순서: 제목(위) → 본문 → 표(아래)
    assert blocks[0]["text"] == "금 시세 조회 자동화"
    table = next(b for b in blocks if b["type"] == "table")
    assert table["rows"][1] == ["시세 조회", "Edge"]

    assert "네이버 증권" in result["full_text"]
    assert "시세 조회 | Edge" in result["full_text"]


def test_pptx_extracts_text_inside_group_shapes():
    """그룹으로 묶인 도형 내부의 텍스트도 추출된다."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃
    group = slide.shapes.add_group_shape()
    box = group.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    box.text_frame.text = "그룹 안의 업무 설명"
    buf = io.BytesIO()
    prs.save(buf)

    result = parse_document("grouped.pptx", buf.getvalue())
    assert "그룹 안의 업무 설명" in result["full_text"]


def test_pdf_blank_page_warns_image_possibility():
    result = parse_document("blank.pdf", _make_blank_pdf())

    assert result["parser"] == "pypdf"
    assert result["page_count"] == 1
    assert result["pages"][0]["blocks"] == []
    assert any("이미지 페이지 가능성" in w for w in result["warnings"])
