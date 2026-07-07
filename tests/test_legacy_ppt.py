"""레거시 PPT(.ppt) 지원 테스트 (RPA-44).

실제 .ppt(OLE) 바이너리·soffice가 없는 환경이라, olefile과 변환 단계를 목킹해
검증 로직·라우팅·graceful degradation을 확인한다.
"""

import io
from types import SimpleNamespace

import pytest
from fastapi import HTTPException

import app.services.parser.ppt as ppt_mod
import app.services.upload_security as sec
from app.services.parser import parse_document
from app.services.upload_security import validate_upload

_OLE = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"


def _fake_olefile(stream_names):
    """listdir이 지정한 스트림을 돌려주는 가짜 olefile 모듈."""
    class _Ole:
        def __init__(self, *a, **k):
            pass

        def listdir(self):
            return [n.split("/") for n in stream_names]

        def close(self):
            pass

    return SimpleNamespace(isOleFile=lambda *a, **k: True, OleFileIO=_Ole)


# --- 업로드 검증 ---

def test_ppt_wrong_magic_rejected():
    with pytest.raises(HTTPException) as e:
        validate_upload("deck.ppt", b"PK\x03\x04 not ole", 20)
    assert e.value.detail["code"] == "FILE_TYPE_MISMATCH"


def test_ppt_valid_powerpoint_stream_ok(monkeypatch):
    import sys

    monkeypatch.setitem(
        sys.modules, "olefile", _fake_olefile(["PowerPoint Document", "SummaryInformation"])
    )
    ct = validate_upload("deck.ppt", _OLE + b"rest", 20)
    assert ct == "application/vnd.ms-powerpoint"


def test_ppt_disguised_doc_rejected(monkeypatch):
    import sys

    # PowerPoint 스트림이 없는 OLE(.doc/.xls 위장) → FILE_TYPE_MISMATCH
    monkeypatch.setitem(sys.modules, "olefile", _fake_olefile(["WordDocument", "1Table"]))
    with pytest.raises(HTTPException) as e:
        validate_upload("fake.ppt", _OLE + b"rest", 20)
    assert e.value.detail["code"] == "FILE_TYPE_MISMATCH"


def test_ppt_with_macro_rejected(monkeypatch):
    import sys

    monkeypatch.setitem(
        sys.modules, "olefile", _fake_olefile(["PowerPoint Document", "_VBA_PROJECT_CUR/VBA/Module1"])
    )
    with pytest.raises(HTTPException) as e:
        validate_upload("macro.ppt", _OLE + b"rest", 20)
    assert e.value.detail["code"] == "SUSPICIOUS_FILE"


# --- 파서 라우팅 + graceful degradation ---

def test_parse_ppt_without_soffice_raises_clear_message(monkeypatch):
    monkeypatch.setattr(ppt_mod, "_find_soffice", lambda: None)
    with pytest.raises(RuntimeError) as e:
        parse_document("old.ppt", _OLE + b"rest")
    assert "LibreOffice" in str(e.value)


def test_ppt_malformed_ole_maps_to_400(monkeypatch):
    """OleFileIO/listdir이 손상 파일에서 던지는 예외는 500이 아니라 CORRUPTED_FILE(400)."""
    import sys

    def _raising_olefileio(*a, **k):
        raise ValueError("broken ole")

    fake = SimpleNamespace(isOleFile=lambda *a, **k: True, OleFileIO=_raising_olefileio)
    monkeypatch.setitem(sys.modules, "olefile", fake)
    with pytest.raises(HTTPException) as e:
        validate_upload("broken.ppt", _OLE + b"rest", 20)
    assert e.value.status_code == 400
    assert e.value.detail["code"] == "CORRUPTED_FILE"


def test_convert_isolates_libreoffice_profile(monkeypatch):
    """동시 변환 충돌 방지 — soffice 호출에 -env:UserInstallation 프로필 격리가 실려야 한다."""
    import os

    from pptx import Presentation

    buf = io.BytesIO()
    Presentation().save(buf)
    pptx_bytes = buf.getvalue()

    captured = {}

    def fake_run(cmd, **kw):
        captured["cmd"] = cmd
        outdir = cmd[cmd.index("--outdir") + 1]
        with open(os.path.join(outdir, "in.pptx"), "wb") as f:
            f.write(pptx_bytes)
        return SimpleNamespace(returncode=0)

    monkeypatch.setattr(ppt_mod, "_find_soffice", lambda: "soffice")
    monkeypatch.setattr(ppt_mod.subprocess, "run", fake_run)

    ppt_mod._convert_ppt_to_pptx(_OLE + b"rest")
    assert any(str(a).startswith("-env:UserInstallation=") for a in captured["cmd"])


def test_parse_ppt_converts_then_parses_pptx(monkeypatch):
    # soffice 변환을 목킹: 실제 pptx 바이트를 돌려주면 parse_pptx가 처리
    import io as _io

    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "레거시에서 변환됨"
    buf = _io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    monkeypatch.setattr(ppt_mod, "_convert_ppt_to_pptx", lambda content: pptx_bytes)
    result = parse_document("old.ppt", _OLE + b"rest")
    assert result["parser"] == "libreoffice+python-pptx"
    assert "레거시에서 변환됨" in result["full_text"]
