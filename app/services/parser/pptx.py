"""PPTX 텍스트·표 추출 — 슬라이드별로 도형을 좌상단 좌표순으로 정렬해 읽기 순서를 보존한다."""

import io
import logging

from pptx import Presentation

logger = logging.getLogger(__name__)


def parse_pptx(content: bytes) -> dict:
    prs = Presentation(io.BytesIO(content))
    pages = []
    warnings = []
    full_parts: list[str] = []

    for i, slide in enumerate(prs.slides, start=1):
        positioned: list[tuple[int, int, dict]] = []
        for shape in slide.shapes:
            top = int(shape.top or 0)
            left = int(shape.left or 0)
            if getattr(shape, "has_table", False):
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                positioned.append((top, left, {"type": "table", "rows": rows}))
            elif getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    positioned.append((top, left, {"type": "text", "text": text}))

        positioned.sort(key=lambda item: (item[0], item[1]))
        blocks = [block for _, _, block in positioned]

        # 발표자 노트에 업무 설명이 담긴 경우가 많아 함께 추출
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                blocks.append({"type": "notes", "text": notes})

        if not blocks:
            warnings.append(f"{i}번 슬라이드에서 텍스트를 찾지 못함 (이미지 슬라이드 가능성 — OCR은 후속 지원)")

        pages.append({"page": i, "blocks": blocks})
        for block in blocks:
            if block["type"] == "table":
                full_parts.append("\n".join(" | ".join(row) for row in block["rows"]))
            else:
                full_parts.append(block["text"])

    return {
        "parser": "python-pptx",
        "page_count": len(pages),
        "pages": pages,
        "full_text": "\n\n".join(full_parts),
        "warnings": warnings,
    }
