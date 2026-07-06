"""비전 LLM 기반 이미지 페이지 보강 파싱 (FR-03).

배경: 업무정의서가 화면 캡처 중심이라 텍스트 추출만으로는 정보 대부분이
유실된다 (샘플 실측: 6페이지에서 511자). 텍스트가 임계값 미만인 페이지를
이미지로 렌더링해 비전 LLM에 구조화 추출을 맡기고, 결과를 parsed_content에
"vision_text" 블록으로 병합한다.

비용 가드: 대상 페이지 임계값(VISION_MIN_TEXT_CHARS)·페이지 수 상한
(VISION_MAX_PAGES)·이미지 폭 다운스케일. 호출은 core.llm 경유라 토큰/비용이
llm_usage에 자동 기록된다.
"""

import base64
import copy
import io
import logging
import os
import re
import uuid
from concurrent.futures import ThreadPoolExecutor, as_completed

logger = logging.getLogger(__name__)

_PROMPT = """이 이미지는 RPA 자동화 대상 업무를 설명하는 업무정의서의 한 페이지입니다.
이미지에 보이는 모든 정보를 빠짐없이 텍스트로 추출하세요:

1. 제목·본문·라벨 등 모든 텍스트 (읽기 순서대로)
2. 표는 행 단위로: 셀1 | 셀2 | 셀3
3. 화면 캡처가 있으면: 어떤 시스템/화면인지, 사용자가 무엇을 하는 장면인지 서술
4. 순서도·화살표가 있으면 흐름을 "A → B → C" 형태로

설명이나 해석을 덧붙이지 말고 이미지의 내용만 충실히 추출하세요."""

_MAX_IMAGE_WIDTH = 1400


def _page_text_chars(page: dict) -> int:
    """페이지의 실질 텍스트 양 — 공백 제외 (layout 모드 패딩이 판정을 왜곡하지 않도록)."""
    total = 0
    for block in page.get("blocks", []):
        if "text" in block:
            total += len(re.sub(r"\s", "", block["text"]))
        if "rows" in block:
            total += sum(len(re.sub(r"\s", "", cell)) for row in block["rows"] for cell in row)
    return total


def _pdf_pages_with_images(content: bytes) -> set[int]:
    """이미지 객체가 있는 PDF 페이지 번호 집합. 판단 불가 시 보수적으로 포함."""
    from pypdf import PdfReader

    result: set[int] = set()
    try:
        reader = PdfReader(io.BytesIO(content))
        for i, page in enumerate(reader.pages, start=1):
            try:
                if len(list(page.images)) > 0:
                    result.add(i)
            except Exception:  # noqa: BLE001
                result.add(i)
    except Exception:  # noqa: BLE001 — 문서 전체 판독 실패 시 필터를 비활성화
        return set(range(1, 10_000))
    return result


def pages_needing_vision(parsed: dict) -> list[int]:
    """텍스트가 임계값 미만인 페이지 번호 목록 (이미 vision 블록이 있으면 제외됨)."""
    threshold = int(os.getenv("VISION_MIN_TEXT_CHARS", "200"))
    return [p["page"] for p in parsed.get("pages", []) if _page_text_chars(p) < threshold]


def render_pdf_pages(content: bytes, page_numbers: list[int]) -> dict[int, list[bytes]]:
    """PDF의 지정 페이지들을 PNG로 렌더링한다."""
    import pypdfium2 as pdfium

    pdf = pdfium.PdfDocument(content)
    try:
        images: dict[int, list[bytes]] = {}
        for n in page_numbers:
            pil = pdf[n - 1].render(scale=2.0).to_pil()
            if pil.width > _MAX_IMAGE_WIDTH:
                ratio = _MAX_IMAGE_WIDTH / pil.width
                pil = pil.resize((_MAX_IMAGE_WIDTH, int(pil.height * ratio)))
            buf = io.BytesIO()
            pil.save(buf, format="PNG")
            images[n] = [buf.getvalue()]
        return images
    finally:
        pdf.close()


def extract_pptx_images(content: bytes, page_numbers: list[int]) -> dict[int, list[bytes]]:
    """PPTX 슬라이드에 내장된 그림(캡처 등)을 큰 순서대로 최대 3장 추출한다 (그룹 내부 포함)."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from app.services.parser.pptx import iter_shapes

    prs = Presentation(io.BytesIO(content))
    images: dict[int, list[bytes]] = {}
    for i, slide in enumerate(prs.slides, start=1):
        if i not in page_numbers:
            continue
        pictures = [
            shape for shape in iter_shapes(slide.shapes)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        ]
        pictures.sort(key=lambda s: int(s.width or 0) * int(s.height or 0), reverse=True)
        images[i] = [p.image.blob for p in pictures[:3]]
    return images


def _extract_page(blobs: list[bytes], model: str | None, session_id: uuid.UUID | None) -> str:
    """페이지 이미지들을 비전 LLM에 보내 텍스트를 추출한다 (병렬 워커에서 실행)."""
    from app.core import llm

    content_parts: list[dict] = [{"type": "text", "text": _PROMPT}]
    for blob in blobs:
        b64 = base64.b64encode(blob).decode()
        content_parts.append(
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}}
        )
    return llm.chat(
        [{"role": "user", "content": content_parts}],
        purpose="vision_parse",
        model=model,
        session_id=session_id,
    ).strip()


def enrich_document_stream(
    filename: str,
    file_content: bytes,
    parsed: dict,
    session_id: uuid.UUID | None = None,
):
    """텍스트 부족 페이지를 비전으로 보강하며 ProgressEvent를 순서대로 yield한다.

    페이지당 LLM 호출이라 수십 초가 걸릴 수 있는 작업 — SSE로 진행 상황을 흘린다.
    이벤트 순서: stage(시작) → partial(페이지별 완료)... → done.
    마지막 done 이벤트의 data에 {"parsed": 보강된 parsed_content, "enriched_pages": [...]}.
    """
    from app.core import llm
    from app.schemas import ProgressEvent

    parsed = copy.deepcopy(parsed)  # ORM JSONB 변경 감지를 위해 새 객체로
    max_pages = int(os.getenv("VISION_MAX_PAGES", "15"))
    ext = filename.rsplit(".", 1)[-1].lower()

    targets = pages_needing_vision(parsed)
    if ext == "pdf" and targets:
        # 텍스트가 이미 어느 정도 추출됐고 이미지도 없는 페이지는 비전 가치가 없다
        # (실측: 순수 텍스트 페이지 비전 호출 → 11.8초에 중복 내용 99자뿐)
        force_empty = int(os.getenv("VISION_FORCE_EMPTY_CHARS", "50"))
        with_images = _pdf_pages_with_images(file_content)
        chars_by_page = {p["page"]: _page_text_chars(p) for p in parsed["pages"]}
        targets = [
            n for n in targets if n in with_images or chars_by_page.get(n, 0) < force_empty
        ]
    targets = targets[:max_pages]

    if not targets:
        yield ProgressEvent(
            event="done",
            stage="vision",
            message="보강이 필요한 페이지가 없습니다",
            data={"parsed": parsed, "enriched_pages": []},
        )
        return

    yield ProgressEvent(
        event="stage",
        stage="vision",
        message=f"이미지 중심 페이지 {len(targets)}개를 비전 분석합니다",
        data={"pages": targets},
    )

    if ext == "pdf":
        page_images = render_pdf_pages(file_content, targets)
    elif ext == "pptx":
        page_images = extract_pptx_images(file_content, targets)
    else:
        raise ValueError(f"비전 파싱이 지원하지 않는 형식: .{ext}")

    model = os.getenv("VISION_MODEL", "").strip() or None
    pages_by_no = {p["page"]: p for p in parsed["pages"]}
    enriched: list[int] = []

    # 페이지별 LLM 호출은 서로 독립 → 병렬 실행 (5페이지 기준 ~21초 → ~6초)
    concurrency = max(1, int(os.getenv("VISION_CONCURRENCY", "4")))
    with ThreadPoolExecutor(max_workers=concurrency) as pool:
        futures = {
            pool.submit(_extract_page, page_images[n], model, session_id): n
            for n in targets
            if page_images.get(n)
        }
        for future in as_completed(futures):
            page_no = futures[future]
            try:
                text = future.result()
            except RuntimeError:
                # 키/쿼터 등 구성 오류는 전 페이지가 실패하므로 즉시 중단
                pool.shutdown(cancel_futures=True)
                raise
            except Exception as e:  # noqa: BLE001 — 개별 페이지 실패는 계속 진행
                logger.warning("%s페이지 비전 추출 실패: %s", page_no, e)
                yield ProgressEvent(
                    event="partial",
                    stage="vision",
                    message=f"{page_no}페이지 추출 실패 (다른 페이지는 계속)",
                    data={"page": page_no, "error": True},
                )
                continue
            if text:
                pages_by_no[page_no]["blocks"].append({"type": "vision_text", "text": text})
                enriched.append(page_no)
            yield ProgressEvent(
                event="partial",
                stage="vision",
                message=f"{page_no}페이지 추출 완료",
                data={"page": page_no, "chars": len(text)},
            )
    enriched.sort()

    if enriched:
        parsed["warnings"] = [
            w for w in parsed.get("warnings", [])
            if not any(w.startswith(f"{n}페이지") or w.startswith(f"{n}번") for n in enriched)
        ]
        if "+vision" not in parsed.get("parser", ""):
            parsed["parser"] = parsed.get("parser", "") + "+vision"
        parsed["full_text"] = _rebuild_full_text(parsed)
        parsed["vision"] = {"enriched_pages": enriched}

    yield ProgressEvent(
        event="done",
        stage="vision",
        message=f"{len(enriched)}개 페이지 보강 완료",
        data={"parsed": parsed, "enriched_pages": enriched},
    )


def enrich_document(
    filename: str,
    file_content: bytes,
    parsed: dict,
    session_id: uuid.UUID | None = None,
) -> tuple[dict, dict]:
    """비스트리밍 편의 래퍼 (테스트·배치용) — 스트림을 소진하고 최종 결과만 반환."""
    last = None
    for event in enrich_document_stream(filename, file_content, parsed, session_id):
        last = event
    return last.data["parsed"], {"enriched_pages": last.data["enriched_pages"]}


def _rebuild_full_text(parsed: dict) -> str:
    parts: list[str] = []
    for page in parsed["pages"]:
        for block in page["blocks"]:
            if "rows" in block:
                parts.append("\n".join(" | ".join(row) for row in block["rows"]))
            elif block.get("text"):
                parts.append(block["text"])
    return "\n\n".join(parts)
